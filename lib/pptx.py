# from pptx import Presentation
import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()
    
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Alcohol Tracker CLI"
    subtitle.text = "A Python Command Line Application for Responsible Drinking"
    
    # Slide 2
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Key Features"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    features = [
        "Track alcohol consumption by type and volume",
        "Set personal limits based on gender/weight",
        "Visual indicators (✅/❌) for consumed drinks",
        "BAC (Blood Alcohol Content) estimation",
        "Warning system for excessive consumption",
        "Unit conversion (oz/ml/l)"
    ]
    
    for feature in features:
        p = tf.add_paragraph()
        p.text = feature
        p.level = 0
        p.space_after = Pt(12)
    

    slide = prs.slides.add_slide(prs.slide_layouts[6])  
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title = title_box.text_frame
    p = title.add_paragraph()
    p.text = "Adding Drinks"
    p.font.size = Pt(32)
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    
    
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8)
    pic = slide.shapes.add_picture("screenshots/add_drink.png", left, top, width=width)
#    tuadd picha
    
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title = title_box.text_frame
    p = title.add_paragraph()
    p.text = "Viewing Consumption"
    p.font.size = Pt(32)
    p.font.bold = True
    

    pic = slide.shapes.add_picture("screenshots/view_drinks.png", left, top, width=width)
    # slide 5
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Technology Stack"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    tech = [
        "Python 3.x",
        "Typer (CLI framework)",
        "Rich (terminal formatting)",
        "SQLAlchemy (database)",
        "python-pptx (this presentation)"
    ]
    
    for item in tech:
        p = tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(12)
    
    # Slide 6
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Key Benefits"
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    benefits = [
        "Promotes responsible drinking",
        "Easy-to-use command line interface",
        "Personalized health warnings",
        "Data-driven consumption insights",
        "Privacy-focused (local database)"
    ]
    
    for benefit in benefits:
        p = tf.add_paragraph()
        p.text = benefit
        p.level = 0
        p.space_after = Pt(12)
    

    prs.save("Alcohol_Tracker_CLI_Presentation.pptx")
    print("Presentation created successfully!")

if __name__ == "__main__":
    create_presentation()